import os
import tempfile
from typing import List

import streamlit as st

# LangChain (stable)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_community.llms.huggingface_pipeline import HuggingFacePipeline
from langchain.chains import ConversationalRetrievalChain
from langchain.docstore.document import Document

# HuggingFace
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline

# File readers
import docx2txt
from pptx import Presentation
from PyPDF2 import PdfReader


# -----------------------------
# File loaders
# -----------------------------

def read_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i}]\n{text}")
    return "\n\n".join(pages)


def read_docx(file_path: str) -> str:
    return docx2txt.process(file_path) or ""


def read_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                slide_text.append(shape.text)
        if slide_text:
            texts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
    return "\n\n".join(texts)


def load_documents(uploaded_files) -> List[Document]:
    docs = []

    for uploaded_file in uploaded_files:
        suffix = uploaded_file.name.split(".")[-1].lower()

        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(uploaded_file.read())
            tmp_path = tmp.name

        if suffix == "pdf":
            text = read_pdf(tmp_path)
        elif suffix == "docx":
            text = read_docx(tmp_path)
        elif suffix in ("ppt", "pptx"):
            text = read_pptx(tmp_path)
        else:
            text = ""

        if text.strip():
            docs.append(
                Document(
                    page_content=text,
                    metadata={"source": uploaded_file.name},
                )
            )

        os.remove(tmp_path)

    return docs


# -----------------------------
# RAG Components (NO CACHING HERE)
# -----------------------------

def build_retriever(docs: List[Document]):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100,
        separators=["\n\n", "\n", ". ", ".", " "],
    )

    chunks = []
    for doc in docs:
        for chunk in splitter.split_text(doc.page_content):
            chunks.append(Document(page_content=chunk, metadata=doc.metadata))

    embeddings = SentenceTransformerEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    vectordb = FAISS.from_documents(chunks, embeddings)
    return vectordb.as_retriever(search_kwargs={"k": 2})


@st.cache_resource
def build_llm():
    model_name = "google/flan-t5-base"

    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name)

    pipe = pipeline(
        "text2text-generation",
        model=model,
        tokenizer=tokenizer,
        max_new_tokens=256,
        do_sample=False,
    )

    return HuggingFacePipeline(pipeline=pipe)


def build_rag_chain(retriever, llm):
    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        return_source_documents=True,
    )


# -----------------------------
# Streamlit UI
# -----------------------------

def main():
    st.set_page_config(page_title="📚 RAG Study Bot", layout="wide")
    st.title("📚 RAG Study Bot (Local, No OpenAI)")

    st.sidebar.header("📄 Upload Documents")
    uploaded_files = st.sidebar.file_uploader(
        "Upload PDF / DOCX / PPTX",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
    )

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    if uploaded_files:
        docs = load_documents(uploaded_files)

        if not docs:
            st.error("No readable content found.")
            return

        retriever = build_retriever(docs)
        llm = build_llm()
        rag_chain = build_rag_chain(retriever, llm)

        st.success("Documents indexed. Ask your questions below 👇")

        query = st.text_input("❓ Ask a question")

        if query:
            result = rag_chain(
                {"question": query, "chat_history": st.session_state.chat_history}
            )

            answer = result["answer"]
            sources = result.get("source_documents", [])

            st.markdown("### 🤖 Answer")
            st.write(answer)

            if sources:
                st.markdown("### 📎 Sources")
                for i, doc in enumerate(sources, start=1):
                    st.markdown(
                        f"**{i}. {doc.metadata.get('source','')}**  \n"
                        f"{doc.page_content[:300]}..."
                    )

            st.session_state.chat_history.append(("user", query))
            st.session_state.chat_history.append(("assistant", answer))


if __name__ == "__main__":
    main()
